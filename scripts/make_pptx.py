"""Generate professional submission presentation for Water Segmentation project."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TEAL = RGBColor(0x0D, 0x94, 0x8B)
TEAL_LIGHT = RGBColor(0xF0, 0xFD, 0xF4)
DARK = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
GRAY_LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
ROW_ALT = RGBColor(0xF0, 0xFD, 0xF4)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x77, 0x06)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
IMG_DIR = "notebooks"


def _dark_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK


def _light_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = GRAY_LIGHT


def _white_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE


def _rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _subtitle_header(slide, number, title):
    _rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), WHITE)
    line = _rect(slide, Inches(0.6), Inches(0.85), Inches(12.133), Inches(0.02), BORDER)
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number}. {title}"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK
    p.font.bold = True
    tag = slide.shapes.add_textbox(Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.4))
    tp = tag.text_frame.paragraphs[0]
    tp.text = f"Slide {number:02d}"
    tp.font.size = Pt(10)
    tp.font.color.rgb = TEAL
    tp.font.bold = True
    tp.alignment = PP_ALIGN.RIGHT
    return tf


def _bullets(slide, items, left, top, width, size=Pt(14), color=None):
    if color is None:
        color = GRAY
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = size
        p.font.color.rgb = color
        p.space_after = Pt(8)
        if item.startswith("  "):
            p.level = 1
            p.font.size = Pt(12)
    return tf


def _section_title(slide, text, left, top, width=Inches(5)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK
    return txBox


def _info_box(slide, text, left, top, width, bg=TEAL_LIGHT):
    r = _rect(slide, left, top, width, Inches(0.6), bg)
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), Inches(0.45))
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = TEAL if bg == TEAL_LIGHT else DARK
    return txBox


def _make_table(slide, headers, rows_data, left, top, width, height):
    rows = len(rows_data) + 1
    cols = len(headers)
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = DARK
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRAY_LIGHT
    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_ALT
    return tbl


def _add_image(slide, path, left, top, width=None, height=None):
    kwargs = {"left": left, "top": top}
    if width:
        kwargs["width"] = width
    if height:
        kwargs["height"] = height
    slide.shapes.add_picture(path, **kwargs)


def _pill(slide, text, left, top, bg, width=Inches(2.8)):
    r = _rect(slide, left, top, width, Inches(0.35), bg)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    return txBox


# ===== SLIDE 1: Title =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
_dark_bg(slide)
_rect(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.03), TEAL)
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(1))
p = txBox.text_frame.paragraphs[0]
p.text = "Water Body Segmentation"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(0.6))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "Binary Semantic Segmentation of Sentinel-2 Satellite Imagery  |  U-Net + ResNet34"
p2.font.size = Pt(16)
p2.font.color.rgb = TEAL
p2.alignment = PP_ALIGN.CENTER

_rect(slide, Inches(0), Inches(5.5), Inches(13.333), Inches(2), DARK)
info_items = [
    ("Candidate", "Angel Gupta"),
    ("Backbone", "U-Net + ResNet34"),
    ("Repository", "Angelgupta13/water-segmentation"),
]
for i, (label, val) in enumerate(info_items):
    left = Inches(1.5 + i * 4)
    txL = slide.shapes.add_textbox(left, Inches(5.8), Inches(3), Inches(0.3))
    pL = txL.text_frame.paragraphs[0]
    pL.text = label
    pL.font.size = Pt(9)
    pL.font.color.rgb = GRAY
    txV = slide.shapes.add_textbox(left, Inches(6.15), Inches(3), Inches(0.3))
    pV = txV.text_frame.paragraphs[0]
    pV.text = val
    pV.font.size = Pt(14)
    pV.font.bold = True
    pV.font.color.rgb = WHITE


# ===== SLIDE 2: Problem & Dataset =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(slide)
_subtitle_header(slide, 2, "Problem Statement & Dataset")
_rect(slide, Inches(0.6), Inches(1.1), Inches(5.5), Inches(0.03), TEAL)
_bullets(slide, [
    "Extracting dynamic water boundaries from Sentinel-2 earth observations is essential for automated flood risk mapping, reservoir monitoring, and climate metrics.",
    "Key challenges:",
    "  Class imbalance: 7.6% water vs 92.4% non-water pixels",
    "  JPEG artifacts: 96% of masks have compression noise (values 1-199)",
    "  RGB-only input -- no NIR band to distinguish water from shadow",
    "Data split: 70% train (1,989), 15% val (426), 15% test (426) from 2,841 images",
], Inches(0.6), Inches(1.3), Inches(5.8), Pt(13))
_add_image(slide, f"{IMG_DIR}/sample_pairs.png", Inches(7), Inches(1.3), height=Inches(5.2))
_info_box(slide, "Data Split: 70% Train (1,989) | 15% Val (426) | 15% Test (426)", Inches(0.6), Inches(6.5), Inches(5.8))


# ===== SLIDE 3: Preprocessing =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(slide)
_subtitle_header(slide, 3, "Preprocessing & Artifact Scrubbing")
_rect(slide, Inches(0.6), Inches(1.1), Inches(5.5), Inches(0.03), TEAL)
_bullets(slide, [
    "Original mask files are riddled with compression artifacts (values 1-199). Standard thresholding limits detail, requiring an optimized configuration:",
    "Resize: 256x256 (bilinear for images, nearest-neighbor for masks)",
    "Threshold: Masks thresholded at >200 to remove JPEG noise",
    "Augmentation: Horizontal/vertical flips (50%), 90-degree rotations, color jitter (30%)",
    "Normalization: ImageNet mean (0.485, 0.456, 0.406) and std (0.229, 0.224, 0.225)",
], Inches(0.6), Inches(1.3), Inches(5.8), Pt(13))
_add_image(slide, f"{IMG_DIR}/mask_artifact.png", Inches(7), Inches(1.3), width=Inches(5.5))
_info_box(slide, "ImageNet Normalization: mean=[0.485, 0.456, 0.406], std=[0.229, 0.224, 0.225]", Inches(0.6), Inches(6.5), Inches(5.8))


# ===== SLIDE 4: EDA =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(slide)
_subtitle_header(slide, 4, "Exploratory Data Analysis")
_rect(slide, Inches(0.6), Inches(1.1), Inches(5.5), Inches(0.03), TEAL)
_bullets(slide, [
    "A complete histogram scan over the water-to-land ratio demonstrates extreme pixel imbalance across scenes:",
    "Water Pixel Ratio: Heavily left-skewed, majority of scenes consist mainly of land background features.",
    "Class Imbalance: Water pixels = 7.6% of total dataset, non-water = 92.4%",
    "Imbalance Handling: Hybrid loss (BCE + Dice) used rather than standard cross-entropy to protect spatial boundaries.",
    "Spatial diversity: Images cover coastlines, rivers, lakes, reservoirs across multiple geographic regions.",
], Inches(0.6), Inches(1.3), Inches(5.8), Pt(13))
_add_image(slide, f"{IMG_DIR}/class_imbalance.png", Inches(7), Inches(1.3), width=Inches(5.5))
_info_box(slide, "Water pixels: ~7.6% of total | BCE + Dice hybrid loss handles imbalance without focal loss", Inches(0.6), Inches(6.5), Inches(5.8))


# ===== SLIDE 5: Model Architecture =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(slide)
_subtitle_header(slide, 5, "Model Architecture")
_rect(slide, Inches(0.6), Inches(1.1), Inches(5.5), Inches(0.03), TEAL)
_bullets(slide, [
    "Encoder: ImageNet-pretrained ResNet34 (21M parameters) ensures fast, stable convergence.",
    "Decoder: Symmetric up-sampling with spatial skip connections from encoder stages.",
    "Loss: 0.5 x BCE + 0.5 x Dice Loss to combat the 7.6% class imbalance.",
    "Optimizer: AdamW (lr=1e-4, weight_decay=1e-4) + CosineAnnealingLR scheduler.",
    "Early stopping: Patience=10 epochs on val IOU (stopped at epoch 34).",
], Inches(0.6), Inches(1.3), Inches(5.8), Pt(13))

# Hyperparameter card on right
_section_title(slide, "Hyperparameter Configuration", Inches(7.5), Inches(1.3))
params = [
    ("Learning Rate", "1e-4"),
    ("Batch Size", "8"),
    ("Optimizer", "AdamW"),
    ("Scheduler", "CosineAnnealingLR"),
    ("Patience", "10 epochs"),
    ("Stopped At", "Epoch 34"),
]
for i, (k, v) in enumerate(params):
    y = Inches(1.9 + i * 0.55)
    _rect(slide, Inches(7.5), y, Inches(5), Inches(0.45), WHITE)
    _rect(slide, Inches(7.5), y, Inches(0.06), Inches(0.45), TEAL)
    tk = slide.shapes.add_textbox(Inches(7.8), y + Inches(0.1), Inches(2.5), Inches(0.3))
    pk = tk.text_frame.paragraphs[0]
    pk.text = k
    pk.font.size = Pt(11)
    pk.font.color.rgb = GRAY
    tv = slide.shapes.add_textbox(Inches(10.3), y + Inches(0.1), Inches(2), Inches(0.3))
    pv = tv.text_frame.paragraphs[0]
    pv.text = v
    pv.font.size = Pt(11)
    pv.font.bold = True
    pv.font.color.rgb = TEAL
    pv.alignment = PP_ALIGN.RIGHT

_info_box(slide, "Weighted BCE handles local spatial gradients; Dice directly addresses global Jaccard metric boundaries", Inches(0.6), Inches(6.5), Inches(5.8))


# ===== SLIDE 6: Results =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(slide)
_subtitle_header(slide, 6, "Hyperparameter Tuning & Results")
_rect(slide, Inches(0.6), Inches(1.1), Inches(5.5), Inches(0.03), TEAL)

_make_table(slide,
    ["Run", "LR", "Batch", "IOU", "Accuracy"],
    [["lr1e-4_bs8", "1e-4", "8", "0.8018", "92.65%"],
     ["lr5e-5_bs8", "5e-5", "8", "0.7928", "92.31%"],
     ["lr1e-4_bs16", "1e-4", "16", "0.7542", "90.86%"]],
    Inches(0.6), Inches(1.3), Inches(5.8), Inches(1.8))

_section_title(slide, "Final Model (Test Set)", Inches(0.6), Inches(3.3))
_make_table(slide,
    ["IOU", "Accuracy", "Precision", "Recall"],
    [["0.8249", "93.80%", "91.90%", "88.70%"]],
    Inches(0.6), Inches(3.8), Inches(5.8), Inches(0.8))

_info_box(slide, "Precision > Recall (91.90% vs 88.70%): under-segments slightly -- safer bias for flood mapping (fewer false alarms)",
          Inches(0.6), Inches(4.9), Inches(5.8), TEAL_LIGHT)

_add_image(slide, f"{IMG_DIR}/predictions.png", Inches(7), Inches(1.3), width=Inches(5.5))


# ===== SLIDE 7: Inference Optimization =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(slide)
_subtitle_header(slide, 7, "Optimized Inference Service")
_rect(slide, Inches(0.6), Inches(1.1), Inches(5.5), Inches(0.03), TEAL)
_bullets(slide, [
    "ONNX Runtime delivers ~2.2x CPU speedup through operator fusion (conv+bn+relu -> single kernel).",
    "Sliding-window tiling: 256x256 tiles, 32px overlap (~72 tiles/image).",
    "Without tiling, a 2009x2007 image would require 46 GB GPU memory.",
    "Test-Time Augmentation: horizontal flip averaging (+0.5-1% IOU).",
    "Auto-fallback: ONNX -> PyTorch -> MLflow registry.",
], Inches(0.6), Inches(1.3), Inches(5.8), Pt(13))

# Latency comparison bars on right
_section_title(slide, "Latency per 256x256 Tile", Inches(7), Inches(1.3))
bars = [
    ("PyTorch (CPU)", "~120 ms", "1.0x", 1.0, GRAY),
    ("ONNX Runtime (CPU)", "~55 ms", "2.2x", 0.46, TEAL),
    ("ONNX Runtime (GPU)", "~15 ms", "8.0x", 0.125, INDIGO),
]
for i, (label, lat, speed, pct, color) in enumerate(bars):
    y = Inches(1.9 + i * 1.0)
    tl = slide.shapes.add_textbox(Inches(7), y, Inches(4.5), Inches(0.25))
    pl = tl.text_frame.paragraphs[0]
    pl.text = f"{label}  ({lat}, {speed})"
    pl.font.size = Pt(11)
    pl.font.color.rgb = DARK
    pl.font.bold = True
    bar_bg = _rect(slide, Inches(7), y + Inches(0.35), Inches(5), Inches(0.25), BORDER)
    bar_fill = _rect(slide, Inches(7), y + Inches(0.35), Inches(5 * pct), Inches(0.25), color)
    tp = slide.shapes.add_textbox(Inches(7 + 5 * pct + 0.1), y + Inches(0.3), Inches(1), Inches(0.3))
    pp = tp.text_frame.paragraphs[0]
    pp.text = speed
    pp.font.size = Pt(10)
    pp.font.bold = True
    pp.font.color.rgb = color

_info_box(slide, "Fallback chain: ONNX -> PyTorch -> MLflow Registry Weight Pull", Inches(0.6), Inches(6.5), Inches(5.8))


# ===== SLIDE 8: Failure Analysis =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(slide)
_subtitle_header(slide, 8, "Failure Modes & Model Transparency")
_rect(slide, Inches(0.6), Inches(1.1), Inches(5.5), Inches(0.03), TEAL)
_make_table(slide,
    ["Failure Mode", "Impact", "Root Cause"],
    [["Narrow waterways", "Streams <4px wide missed", "256x256 resize destroys thin features"],
     ["Shadow / dark terrain", "Shadows classified as water", "RGB-only; NIR band disambiguates"],
     ["JPEG boundary noise", "Ragged mask edges (+/-1px)", "Compression artifacts in source masks"]],
    Inches(0.6), Inches(1.3), Inches(5.8), Inches(1.8))
_section_title(slide, "Operational Boundaries", Inches(0.6), Inches(3.3))
_bullets(slide, [
    "Production-ready for: lakes, large rivers, wide reservoirs",
    "Not suitable for: narrow canals, small ponds (<500 m2), shadow-prone terrain",
    "Biggest bottleneck: RGB-only input -- adding NIR band gives +5-10% IOU",
], Inches(0.6), Inches(3.8), Inches(5.8), Pt(13))

_add_image(slide, f"{IMG_DIR}/edge_cases.png", Inches(7), Inches(1.3), width=Inches(5.5))


# ===== SLIDE 9: MLOps =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(slide)
_subtitle_header(slide, 9, "MLOps Deployment Pipeline")
_rect(slide, Inches(0.6), Inches(1.1), Inches(5.5), Inches(0.03), TEAL)

# Pipeline flow
steps = ["COMMIT", "LINT", "TEST", "BUILD", "PUSH"]
step_colors = [TEAL, TEAL, TEAL, TEAL, TEAL]
for i, (step, sc) in enumerate(zip(steps, step_colors)):
    x = Inches(0.6 + i * 2.3)
    r = _rect(slide, x, Inches(1.3), Inches(2), Inches(0.9), WHITE)
    r.line.color.rgb = BORDER
    _rect(slide, x, Inches(1.3), Inches(2), Inches(0.04), sc)
    ts = slide.shapes.add_textbox(x, Inches(1.4), Inches(2), Inches(0.35))
    ps = ts.text_frame.paragraphs[0]
    ps.text = step
    ps.font.size = Pt(12)
    ps.font.bold = True
    ps.font.color.rgb = sc
    ps.alignment = PP_ALIGN.CENTER

_bullets(slide, [
    "GitHub Actions on push to master: flake8 lint -> 15 pytest -> Docker build -> Docker Hub push.",
    "Dockerfile: python:3.11-slim, ONNX model baked in, HEALTHCHECK every 10s.",
    "MLflow: per-epoch metrics, model registry (water-segmentation-unet v1, serialization_format=pt2).",
    "Security: 60 req/min rate limiter, 100 MB upload cap, path traversal protection.",
    "Image: angelgupta/water-segmentation:latest on Docker Hub.",
], Inches(0.6), Inches(2.5), Inches(5.8), Pt(13))
_add_image(slide, f"{IMG_DIR}/demo_output.png", Inches(7), Inches(1.3), width=Inches(5.5))


# ===== SLIDE 10: Recommendations =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
_light_bg(slide)
_subtitle_header(slide, 10, "Recommendations & Strategic Roadmap")

cols = [
    (Inches(0.6), RED, "HIGH IMPACT (+5-10% IOU)", [
        "Add NIR band input (Band 8) for NDWI",
        "Train on full-resolution random crops",
    ]),
    (Inches(4.8), AMBER, "MEDIUM IMPACT (+1-3% IOU)", [
        "EfficientNet-B4 encoder swap",
        "Optuna Bayesian search (30+ trials)",
    ]),
    (Inches(9), TEAL, "LOW IMPACT (+0.5-1% IOU)", [
        "Ensemble (3-5 seeds)",
        "CRF post-processing",
    ]),
]
for x, color, title, items in cols:
    _rect(slide, x, Inches(1.3), Inches(3.8), Inches(0.5), color)
    tt = slide.shapes.add_textbox(x, Inches(1.3), Inches(3.8), Inches(0.5))
    pt = tt.text_frame.paragraphs[0]
    pt.text = title
    pt.font.size = Pt(11)
    pt.font.bold = True
    pt.font.color.rgb = WHITE
    pt.alignment = PP_ALIGN.CENTER
    _bullets(slide, [f"  {item}" for item in items],
             x + Inches(0.2), Inches(2.0), Inches(3.4), Pt(13))

txC = slide.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(12), Inches(0.5))
pC = txC.text_frame.paragraphs[0]
pC.text = "Current model (IOU 0.82) is production-ready for coarse water mapping. Adding NIR band closes the gap to high-precision applications."
pC.font.size = Pt(14)
pC.font.color.rgb = DARK
pC.alignment = PP_ALIGN.CENTER

_info_box(slide, "7 prioritized recommendations with expected impact and effort in REPORT.md Section 10",
          Inches(0.6), Inches(6.5), Inches(12), TEAL_LIGHT)


# Save
import sys
out_path = sys.argv[1] if len(sys.argv) > 1 else "presentation.pptx"
prs.save(out_path)
print(f"Presentation saved to {out_path}")
